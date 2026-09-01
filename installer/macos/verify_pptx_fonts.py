#!/usr/bin/env python3
"""Render a Korean font smoke-test slide through the installed pptx-mcp code."""

from __future__ import annotations

import argparse
import json
import shutil
import sys
import tempfile
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


FONT_SAMPLES = [
    ("Noto Sans CJK KR", "한글 렌더링 검증 하나"),
    ("Noto Serif CJK KR", "한글 렌더링 검증 둘"),
    ("Nanum Gothic", "한글 렌더링 검증 셋"),
    ("Nanum Myeongjo", "한글 렌더링 검증 넷"),
    ("NanumSquare", "한글 렌더링 검증 다섯"),
    ("NanumBarunGothic", "한글 렌더링 검증 여섯"),
    ("UnBatang", "한글 렌더링 검증 일곱"),
    ("UnDotum", "한글 렌더링 검증 여덟"),
]

PDF_FONT_TOKENS = {
    "Noto Sans CJK KR": "NotoSansCJKkr",
    "Noto Serif CJK KR": "NotoSerifCJKkr",
    "Nanum Gothic": "NanumGothic",
    "Nanum Myeongjo": "NanumMyeongjo",
    "NanumSquare": "NanumSquare",
    "NanumBarunGothic": "NanumBarunGothic",
    "UnBatang": "UnBatang",
    "UnDotum": "UnDotum",
}


def set_run_typeface(run, family: str) -> None:
    run.font.name = family
    properties = run._r.get_or_add_rPr()  # noqa: SLF001 - python-pptx has no public eastAsia API
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = properties.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            properties.append(element)
        element.set("typeface", family)


def build_smoke_presentation(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.5))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = "MnP Suite macOS 한글 글꼴 PNG 검증"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    set_run_typeface(title_run, "Noto Sans CJK KR")

    for index, (family, sample) in enumerate(FONT_SAMPLES):
        text_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.95 + index * 0.72), Inches(11.9), Inches(0.55)
        )
        paragraph = text_box.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = f"{family} — {sample} 123 ABC"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(20, 20, 20)
        set_run_typeface(run, family)

    presentation.save(path)


def verify_png(path: Path) -> tuple[int, int]:
    with Image.open(path) as image:
        image.verify()
    with Image.open(path) as image:
        rgb = image.convert("RGB")
        difference = ImageChops.difference(rgb, Image.new("RGB", rgb.size, "white"))
        if difference.getbbox() is None:
            raise RuntimeError("Rendered PNG is blank")
        return rgb.size


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx-root", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()

    pptx_root = args.pptx_root.resolve()
    output = args.output.resolve()
    sys.path.insert(0, str(pptx_root))

    from tools.export_tools import _convert_pptx_to_pdf, _render_pdf_to_pngs

    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="mnp_pptx_font_smoke_") as temporary_directory:
        temporary_root = Path(temporary_directory)
        pptx_path = temporary_root / "korean-font-smoke.pptx"
        pdf_directory = temporary_root / "pdf"
        png_directory = temporary_root / "png"
        pdf_directory.mkdir()
        png_directory.mkdir()
        build_smoke_presentation(pptx_path)

        pdf_path = Path(_convert_pptx_to_pdf(str(pptx_path), str(pdf_directory)))
        png_paths = _render_pdf_to_pngs(str(pdf_path), str(png_directory), dpi=150)
        if len(png_paths) != 1:
            raise RuntimeError(f"Expected one rendered PNG, got {len(png_paths)}")

        import pymupdf as fitz

        with fitz.open(pdf_path) as document:
            extracted_text = "\n".join(page.get_text() for page in document)
            pdf_fonts = sorted(
                {
                    font[3]
                    for page in document
                    for font in page.get_fonts(full=True)
                    if len(font) > 3 and font[3]
                }
            )
        missing_samples = [sample for _, sample in FONT_SAMPLES if sample not in extracted_text]
        if missing_samples:
            raise RuntimeError(f"Korean text missing after LibreOffice conversion: {missing_samples}")
        missing_pdf_families = [
            family
            for family, token in PDF_FONT_TOKENS.items()
            if not any(token.lower() in pdf_font.lower() for pdf_font in pdf_fonts)
        ]
        if missing_pdf_families:
            raise RuntimeError(
                "LibreOffice substituted or omitted required Korean fonts: "
                f"{missing_pdf_families}; PDF fonts: {pdf_fonts}"
            )

        rendered_png = Path(png_paths[0])
        width, height = verify_png(rendered_png)
        shutil.copy2(rendered_png, output)

    print(
        json.dumps(
            {
                "renderer": "libreoffice+pymupdf",
                "output": str(output),
                "width": width,
                "height": height,
                "verifiedFamilies": [family for family, _ in FONT_SAMPLES],
                "pdfFonts": pdf_fonts,
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
